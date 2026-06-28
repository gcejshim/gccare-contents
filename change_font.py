from pptx import Presentation
from pptx.util import Pt
import io
import os

INPUT_FILE = r"c:\Users\윤성웅\.claude\projects\gccare-liver\docs\GC_Care_AI_Content.pptx"
OUTPUT_FILE = r"c:\Users\윤성웅\.claude\projects\gccare-liver\docs\GC_Care_AI_Content_Pretendard.pptx"
TARGET_FONT = "Pretendard"

with open(INPUT_FILE, 'rb') as f:
    prs = Presentation(io.BytesIO(f.read()))

changed = 0

for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            # 단락 기본 폰트
            if para.runs:
                pass
            for run in para.runs:
                if run.font.name != TARGET_FONT:
                    run.font.name = TARGET_FONT
                    changed += 1
            # 단락 레벨 폰트 (pPr > defRPr)
            pPr = para._p.get_or_add_pPr()
            # shape 레벨 폰트도 처리
        # text_frame 전체 기본 폰트
        tf = shape.text_frame
        txBody = tf._txBody
        lstStyle = txBody.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}lstStyle')
        if lstStyle is not None:
            for defRPr in lstStyle.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}defRPr'):
                defRPr.set('lang', 'ko-KR')
                latin = defRPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}latin')
                if latin is None:
                    import lxml.etree as etree
                    latin = etree.SubElement(defRPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}latin')
                latin.set('typeface', TARGET_FONT)

# 슬라이드 마스터 / 레이아웃 폰트도 변경
for master in prs.slide_masters:
    for shape in master.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = TARGET_FONT

with open(OUTPUT_FILE, 'wb') as f:
    buf = io.BytesIO()
    prs.save(buf)
    f.write(buf.getvalue())
print(f"완료! 변경된 텍스트 런: {changed}개")
print(f"저장 위치: {os.path.abspath(OUTPUT_FILE)}")
