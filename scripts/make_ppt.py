from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── 색상 팔레트 ──────────────────────────────────
GREEN      = RGBColor(0x0A, 0xC2, 0x62)   # #0AC262
NAVY       = RGBColor(0x0D, 0x1F, 0x3C)   # #0D1F3C
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF4, 0xF6, 0xF9)
DARK_GRAY  = RGBColor(0x44, 0x44, 0x44)
MID_GRAY   = RGBColor(0x88, 0x88, 0x88)
GREEN_LIGHT= RGBColor(0xE8, 0xFA, 0xF1)   # 연한 초록 배경

W = Inches(13.33)
H = Inches(7.5)

def prs_new():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    blank = prs.slide_layouts[6]
    return prs.slides.add_slide(blank)

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, l, t, w, h, bg_color=None, border_color=None, border_pt=0):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.line.fill.background()
    if bg_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
    else:
        shape.fill.background()
    if border_color and border_pt:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, l, t, w, h,
        size=18, bold=False, color=None, align=PP_ALIGN.LEFT,
        wrap=True, valign=None):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    if valign:
        tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or DARK_GRAY
    return tb

def txt_multiline(slide, lines, l, t, w, h,
                  size=18, bold=False, color=None, align=PP_ALIGN.LEFT,
                  line_spacing=None):
    """lines: list of (text, bold_override, size_override, color_override)"""
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    from lxml import etree

    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True

    for i, item in enumerate(lines):
        if isinstance(item, str):
            item = (item, bold, size, color)
        text, b, s, c = item[0], item[1] if len(item)>1 else bold, item[2] if len(item)>2 else size, item[3] if len(item)>3 else color

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = text
        run.font.size = Pt(s)
        run.font.bold = b
        run.font.color.rgb = c or DARK_GRAY
    return tb

def add_green_bar(slide, t=Inches(0.55), thickness=Pt(3)):
    bar = slide.shapes.add_shape(1, 0, t, W, thickness)
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN
    bar.line.fill.background()

def section_header_block(slide, number, title):
    """왼쪽 녹색 번호 + 제목 블록"""
    # 번호 원형
    circ = slide.shapes.add_shape(9, Inches(0.4), Inches(0.72), Inches(0.55), Inches(0.55))
    circ.fill.solid(); circ.fill.fore_color.rgb = GREEN
    circ.line.fill.background()
    tf = circ.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    from pptx.enum.text import MSO_ANCHOR
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    r = p.add_run(); r.text = number
    r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = WHITE

    txt(slide, title, Inches(1.05), Inches(0.68), Inches(11), Inches(0.6),
        size=22, bold=True, color=NAVY)


# ══════════════════════════════════════════════════════════════════
# 슬라이드 제작
# ══════════════════════════════════════════════════════════════════

prs = prs_new()

# ──────────────────────────────────────────────
# SLIDE 1 — 표지
# ──────────────────────────────────────────────
s = blank_slide(prs)
bg(s, NAVY)

# 녹색 사각형 왼쪽 엣지
box(s, 0, 0, Inches(0.18), H, bg_color=GREEN)

# 상단 장식 원
circ1 = s.shapes.add_shape(9, Inches(10.8), Inches(-0.8), Inches(2.8), Inches(2.8))
circ1.fill.solid(); circ1.fill.fore_color.rgb = RGBColor(0x14, 0x30, 0x55)
circ1.line.fill.background()

circ2 = s.shapes.add_shape(9, Inches(11.5), Inches(0.6), Inches(1.8), Inches(1.8))
circ2.fill.solid(); circ2.fill.fore_color.rgb = RGBColor(0x0A, 0xC2, 0x62)
circ2.fill.fore_color.rgb = RGBColor(0x1A, 0x3A, 0x62)
circ2.line.fill.background()

# 태그
tag = s.shapes.add_shape(1, Inches(0.6), Inches(1.6), Inches(2.8), Inches(0.42))
tag.fill.solid(); tag.fill.fore_color.rgb = GREEN
tag.line.fill.background()
tf = tag.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
from pptx.enum.text import MSO_ANCHOR
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
r = p.add_run(); r.text = "GC Care · 도입 결과 보고"
r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = WHITE

# 메인 제목
txt_multiline(s, [
    ("AI 자동화 + 웹 편집기 기반", False, 26, RGBColor(0xAA, 0xCC, 0xFF)),
    ("카드뉴스 제작 시스템", True, 42, WHITE),
    ("도입 보고", True, 42, GREEN),
], Inches(0.6), Inches(2.2), Inches(10), Inches(3.2), align=PP_ALIGN.LEFT)

# 부제
txt(s, "AI 이미지 생성 + 웹 편집기 기반 제작 프로세스 도입",
    Inches(0.6), Inches(5.4), Inches(9), Inches(0.6),
    size=16, color=RGBColor(0x88, 0xAA, 0xCC))

# 날짜
txt(s, "2026년 5월 · GC Care 디지털콘텐츠팀",
    Inches(0.6), Inches(6.5), Inches(8), Inches(0.5),
    size=13, color=RGBColor(0x66, 0x88, 0xAA))


# ──────────────────────────────────────────────
# SLIDE 2 — 목차
# ──────────────────────────────────────────────
s = blank_slide(prs)
bg(s, WHITE)
add_green_bar(s, t=Inches(0), thickness=Inches(0.08))

txt(s, "목차", Inches(0.5), Inches(0.2), Inches(3), Inches(0.6),
    size=28, bold=True, color=NAVY)

items = [
    ("01", "추진 배경",          "제작 공수 문제와 도입 필요성"),
    ("02", "시스템 개요",        "기존 방식 vs 신규 방식 비교"),
    ("03", "제작 파이프라인",    "기획서 → 자동화 → 검수 → 편집 → 서비스"),
    ("04", "사용 기술 및 도구",  "AI 모델 선택 이유 및 편집기 구조"),
    ("05", "제작 카드뉴스 현황", "완성 26개 목록"),
    ("06", "비용 분석",          "실비용 및 기존 방식 대비 효율"),
    ("07", "향후 계획",          "추가 자동화 및 개선 방향"),
]

for i, (no, title, desc) in enumerate(items):
    top = Inches(1.0) + i * Inches(0.85)
    # 번호 배경
    nb = s.shapes.add_shape(1, Inches(0.4), top, Inches(0.52), Inches(0.52))
    nb.fill.solid(); nb.fill.fore_color.rgb = GREEN if i % 2 == 0 else NAVY
    nb.line.fill.background()
    tf = nb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    r = p.add_run(); r.text = no
    r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = WHITE

    txt(s, title, Inches(1.1), top - Inches(0.03), Inches(3.5), Inches(0.55),
        size=16, bold=True, color=NAVY)
    txt(s, desc, Inches(4.8), top + Inches(0.03), Inches(8), Inches(0.45),
        size=13, color=MID_GRAY)

    # 구분선
    if i < len(items)-1:
        line = s.shapes.add_shape(1, Inches(0.4), top + Inches(0.6), Inches(12.5), Pt(0.8))
        line.fill.solid(); line.fill.fore_color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
        line.line.fill.background()


# ──────────────────────────────────────────────
# SLIDE 3 — 추진 배경
# ──────────────────────────────────────────────
s = blank_slide(prs)
bg(s, WHITE)
add_green_bar(s, t=Inches(0), thickness=Inches(0.08))
section_header_block(s, "01", "추진 배경")

# 배경 설명 박스
box(s, Inches(0.4), Inches(1.45), Inches(12.5), Inches(1.0), bg_color=LIGHT_GRAY)
txt(s, "간 건강 카드뉴스 26개 콘텐츠 제작 — 기존 디자이너 방식으로는 제작 기간과 반복 수정 비용이 과다하게 발생",
    Inches(0.65), Inches(1.55), Inches(12.0), Inches(0.8),
    size=15, color=DARK_GRAY)

# 수치 카드 3개
cards = [
    ("182개", "전체 디자인 작업 단위", "26개 카드 × 7컷\n(표지 1 + 섹션 6)"),
    ("약 9일", "디자이너 순수 제작 소요일", "1인 하루 20컷 기준\n수정 발생 시 +1~2일 추가"),
    ("텍스트\n수정 多", "수정 요청의 대부분", "디자이너 재작업 반복\n리소스 비효율 구조"),
]
for i, (num, sub, desc) in enumerate(cards):
    left = Inches(0.4) + i * Inches(4.2)
    box(s, left, Inches(2.65), Inches(3.9), Inches(2.8),
        bg_color=NAVY if i == 0 else (GREEN_LIGHT if i == 1 else LIGHT_GRAY))
    num_color = WHITE if i == 0 else (GREEN if i == 1 else NAVY)
    sub_color = RGBColor(0xAA, 0xCC, 0xFF) if i == 0 else MID_GRAY
    desc_color = RGBColor(0xCC, 0xDD, 0xFF) if i == 0 else DARK_GRAY

    txt(s, num, left + Inches(0.2), Inches(2.85), Inches(3.5), Inches(0.9),
        size=34, bold=True, color=num_color, align=PP_ALIGN.CENTER)
    txt(s, sub, left + Inches(0.2), Inches(3.7), Inches(3.5), Inches(0.4),
        size=13, bold=True, color=sub_color, align=PP_ALIGN.CENTER)
    txt(s, desc, left + Inches(0.2), Inches(4.15), Inches(3.5), Inches(0.9),
        size=12, color=desc_color, align=PP_ALIGN.CENTER)

txt(s, "→ 제작 기간 단축 + 담당자 직접 수정이 가능한 구조가 필요",
    Inches(0.4), Inches(5.65), Inches(12.5), Inches(0.6),
    size=15, bold=True, color=GREEN)


# ──────────────────────────────────────────────
# SLIDE 4 — 시스템 개요
# ──────────────────────────────────────────────
s = blank_slide(prs)
bg(s, WHITE)
add_green_bar(s, t=Inches(0), thickness=Inches(0.08))
section_header_block(s, "02", "시스템 개요")

txt(s, "AI 자동화 + 웹 편집기로 제작 효율화",
    Inches(0.4), Inches(1.35), Inches(12), Inches(0.5),
    size=14, color=MID_GRAY)

# 두 축 설명
cols = [
    ("🤖  AI 자동화", GREEN, WHITE,
     ["기획서 데이터 기반 이미지 자동 생성",
      "텍스트·이미지 HTML 자동 주입",
      "파일 저장까지 원클릭 처리"]),
    ("✏️  웹 편집기", NAVY, WHITE,
     ["브라우저에서 바로 텍스트 수정",
      "이미지 교체(붙여넣기/드래그)",
      "JPG 내보내기까지 단일 파일"]),
]
for i, (title, bg_c, tc, bullets) in enumerate(cols):
    left = Inches(0.4) + i * Inches(6.3)
    box(s, left, Inches(1.9), Inches(5.9), Inches(0.62), bg_color=bg_c)
    txt(s, title, left + Inches(0.2), Inches(1.98), Inches(5.5), Inches(0.5),
        size=18, bold=True, color=tc)
    for j, b in enumerate(bullets):
        bt = Inches(2.65) + j * Inches(0.72)
        box(s, left, bt, Inches(5.9), Inches(0.65),
            bg_color=LIGHT_GRAY if j % 2 == 0 else WHITE)
        txt(s, "▸  " + b, left + Inches(0.2), bt + Inches(0.1), Inches(5.5), Inches(0.5),
            size=14, color=DARK_GRAY)

# 비교 표 헤더
table_top = Inches(4.95)
headers = ["구분", "기존 방식", "신규 방식"]
col_w   = [Inches(2.2), Inches(4.8), Inches(5.0)]
col_l   = [Inches(0.4), Inches(2.65), Inches(7.5)]

for i, (h, cw, cl) in enumerate(zip(headers, col_w, col_l)):
    box(s, cl, table_top, cw, Inches(0.42), bg_color=NAVY)
    txt(s, h, cl + Inches(0.1), table_top + Inches(0.05), cw, Inches(0.35),
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rows = [
    ("제작 기간", "약 9~11일", "1일 이내"),
    ("수정 대응", "디자이너 일정 의존", "즉시 가능"),
    ("비용", "디자이너 인건비 발생", "AI API 비용만"),
]
for ri, (c1, c2, c3) in enumerate(rows):
    rt = table_top + Inches(0.42) + ri * Inches(0.42)
    bg_c = LIGHT_GRAY if ri % 2 == 0 else WHITE
    for ci, (cell, cw, cl) in enumerate(zip([c1,c2,c3], col_w, col_l)):
        box(s, cl, rt, cw, Inches(0.42), bg_color=bg_c)
        fc = GREEN if (ci == 2 and ri == 0) else DARK_GRAY
        txt(s, cell, cl + Inches(0.1), rt + Inches(0.07), cw - Inches(0.15), Inches(0.35),
            size=12, color=fc, bold=(ci == 2 and ri == 0))


# ──────────────────────────────────────────────
# SLIDE 5 — 제작 파이프라인
# ──────────────────────────────────────────────
s = blank_slide(prs)
bg(s, WHITE)
add_green_bar(s, t=Inches(0), thickness=Inches(0.08))
section_header_block(s, "03", "제작 파이프라인")

steps = [
    ("STEP 1", "콘텐츠 기획서 작성", "주제·표지 제목·섹션 텍스트 정의"),
    ("STEP 2", "AI 이미지 자동 생성\n→ HTML 자동 주입 → 저장", "기획서 하나로 완성본까지 원클릭"),
    ("STEP 3", "디자이너 검수\n및 수정 보완", "이미지 품질·스타일 확인 후 재생성"),
    ("STEP 4", "담당자 편집기 활용\n텍스트·이미지 수정", "브라우저에서 직접 수정 → JPG 내보내기"),
    ("STEP 5", "서비스 활용", "CMS 등록 → 카카오 메시지 → 링크 페이지"),
]

sw = Inches(2.3)
for i, (step, title, desc) in enumerate(steps):
    left = Inches(0.3) + i * Inches(2.58)
    bg_c = NAVY if i in (1, 3) else (GREEN if i == 4 else LIGHT_GRAY)
    tc   = WHITE if i in (1, 3, 4) else NAVY
    dc   = RGBColor(0xAA, 0xCC, 0xFF) if i == 1 else (RGBColor(0xCC, 0xFF, 0xDD) if i == 4 else MID_GRAY)

    box(s, left, Inches(1.4), sw, Inches(5.5), bg_color=bg_c)
    txt(s, step, left + Inches(0.1), Inches(1.55), sw - Inches(0.1), Inches(0.4),
        size=11, bold=True, color=GREEN if bg_c == NAVY else (WHITE if bg_c == GREEN else GREEN),
        align=PP_ALIGN.CENTER)
    txt(s, title, left + Inches(0.1), Inches(2.1), sw - Inches(0.1), Inches(1.5),
        size=14, bold=True, color=tc, align=PP_ALIGN.CENTER)
    txt(s, desc, left + Inches(0.1), Inches(3.8), sw - Inches(0.1), Inches(1.5),
        size=11, color=dc, align=PP_ALIGN.CENTER)

    # 화살표
    if i < len(steps) - 1:
        ax = left + sw + Inches(0.08)
        txt(s, "▶", ax, Inches(3.8), Inches(0.35), Inches(0.5),
            size=18, color=GREEN, align=PP_ALIGN.CENTER)

# 하단 활용 흐름
box(s, Inches(0.3), Inches(7.0), Inches(12.7), Inches(0.38), bg_color=GREEN_LIGHT)
txt(s, "서비스 흐름:  CMS 이미지 등록  →  카카오 메시지 발송  →  링크 페이지로 고객 서비스 제공",
    Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.32),
    size=12, color=NAVY, align=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────
# SLIDE 6 — 사용 기술 및 도구
# ──────────────────────────────────────────────
s = blank_slide(prs)
bg(s, WHITE)
add_green_bar(s, t=Inches(0), thickness=Inches(0.08))
section_header_block(s, "04", "사용 기술 및 도구")

# AI 모델 비교 박스
txt(s, "AI 모델 구성 및 선택 이유",
    Inches(0.4), Inches(1.38), Inches(8), Inches(0.45),
    size=15, bold=True, color=NAVY)

models = [
    ("표지 이미지", "GPT-image-1\n(OpenAI)", GREEN,
     "한글 텍스트 정확도\n우수 — 자모 오류 없음"),
    ("섹션 이미지", "Gemini 2.5 Flash\n(Google)", NAVY,
     "빠른 생성 속도\n낮은 API 비용"),
]
for i, (role, model, mc, reason) in enumerate(models):
    left = Inches(0.4) + i * Inches(4.5)
    box(s, left, Inches(1.9), Inches(4.1), Inches(0.45), bg_color=mc)
    txt(s, role, left+Inches(0.15), Inches(1.97), Inches(3.8), Inches(0.35),
        size=13, bold=True, color=WHITE)
    box(s, left, Inches(2.35), Inches(4.1), Inches(1.5), bg_color=LIGHT_GRAY)
    txt(s, model, left+Inches(0.15), Inches(2.42), Inches(3.8), Inches(1.3),
        size=18, bold=True, color=mc)
    box(s, left, Inches(3.85), Inches(4.1), Inches(1.1), bg_color=WHITE)
    txt(s, reason, left+Inches(0.15), Inches(3.9), Inches(3.8), Inches(0.95),
        size=13, color=DARK_GRAY)

# 비교 안내 박스
box(s, Inches(0.4), Inches(5.1), Inches(8.6), Inches(1.0),
    bg_color=RGBColor(0xFF, 0xF8, 0xE1), border_color=RGBColor(0xFF, 0xC1, 0x07), border_pt=1)
txt_multiline(s, [
    ("💡 표지 모델 선택 근거", True, 13, RGBColor(0x99, 0x66, 0x00)),
    ("Gemini는 한글 자모 오류·글자 깨짐이 빈번 발생", False, 12, DARK_GRAY),
    ("GPT-image-1은 한글을 정확하게 렌더링 → 표지에 필수", False, 12, DARK_GRAY),
], Inches(0.6), Inches(5.18), Inches(8.2), Inches(0.85))

# 이미지 비교 자리 표시
box(s, Inches(9.3), Inches(1.35), Inches(3.7), Inches(4.2),
    bg_color=LIGHT_GRAY, border_color=RGBColor(0xCC,0xCC,0xCC), border_pt=1)
txt_multiline(s, [
    ("표지 품질 비교", True, 14, NAVY),
    ("", False, 10, MID_GRAY),
    ("[Gemini 생성 결과]", False, 12, MID_GRAY),
    ("한글 오류 발생 예시", False, 11, MID_GRAY),
    ("", False, 10, MID_GRAY),
    ("[GPT-image-1 생성 결과]", False, 12, NAVY),
    ("정확한 한글 렌더링", False, 11, MID_GRAY),
], Inches(9.4), Inches(1.5), Inches(3.5), Inches(3.8), align=PP_ALIGN.CENTER)

# 편집기 특징
txt(s, "편집기 주요 특징",
    Inches(0.4), Inches(6.25), Inches(8), Inches(0.4),
    size=14, bold=True, color=NAVY)
feats = ["브라우저에서 바로 실행 (설치 불필요)",
         "텍스트 직접 수정 · 이미지 교체 지원",
         "JPG 내보내기 · 파일명 자동 생성",
         "단일 HTML 파일로 편집·저장 일체형"]
for i, f in enumerate(feats):
    left = Inches(0.4) + (i % 2) * Inches(6.2)
    top  = Inches(6.7) + (i // 2) * Inches(0.4)
    txt(s, "✓  " + f, left, top, Inches(5.8), Inches(0.38),
        size=13, color=DARK_GRAY)


# ──────────────────────────────────────────────
# SLIDE 7 — 제작 카드뉴스 현황
# ──────────────────────────────────────────────
s = blank_slide(prs)
bg(s, WHITE)
add_green_bar(s, t=Inches(0), thickness=Inches(0.08))
section_header_block(s, "05", "제작 카드뉴스 현황")

txt(s, "총 26개 완성  |  간 건강 관리 시리즈",
    Inches(0.4), Inches(1.35), Inches(10), Inches(0.45),
    size=14, color=MID_GRAY)

cards26 = [
    ("01","간수치가 높다는 신호"),("02","AST·ALT·γ-GTP 쉽게 이해하기"),
    ("03","지방간이란 무엇인가"),("04","내 간수치 상승 원인 찾기"),
    ("05","지방간의 종류와 관리법"),("06","지방간을 방치하면?"),
    ("07","지방간 개선의 핵심 3가지"),("08","체중 5% 감량 효과"),
    ("09","금주·절주 실천 가이드"),("10","하루 7,000보 걷기 챌린지"),
    ("11","탄수화물 줄이기"),("12","단백질 중심 식단 시작"),
    ("13","야식과 지방간"),("14","간에 좋은 식사 패턴"),
    ("15","간 건강과 수면의 관계"),("16","스트레스와 간 건강"),
    ("17","간에 좋은 음식 TOP 10"),("18","간 건강 검진 가이드"),
    ("19","약물과 간 부담 주의"),("20","커피와 간 건강"),
    ("21","간 건강 수분 섭취"),("22","건강기능식품 주의사항"),
    ("23","간경변 예방과 관리"),("24","B형 간염 예방접종"),
    ("25","간 건강 자가 체크리스트"),("26","GC케어 간 건강 프로그램"),
]

cols_n  = 4
cell_w  = Inches(3.18)
cell_h  = Inches(0.48)
gap_h   = Inches(0.05)
margin_l = Inches(0.32)
start_t  = Inches(1.88)

for i, (no, title) in enumerate(cards26):
    col = i % cols_n
    row = i // cols_n
    cl = margin_l + col * (cell_w + Inches(0.1))
    ct = start_t + row * (cell_h + gap_h)
    row_alt = (row % 2 == 0)
    box(s, cl, ct, cell_w, cell_h,
        bg_color=LIGHT_GRAY if row_alt else WHITE,
        border_color=RGBColor(0xDD, 0xDD, 0xDD), border_pt=0.5)
    # 번호를 텍스트로 직접 표시 (shape 없이)
    txt(s, no, cl + Inches(0.08), ct + Inches(0.07),
        Inches(0.38), Inches(0.34),
        size=11, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    txt(s, title, cl + Inches(0.5), ct + Inches(0.09),
        cell_w - Inches(0.58), Inches(0.32),
        size=11, color=DARK_GRAY)


# ──────────────────────────────────────────────
# SLIDE 8 — 비용 분석
# ──────────────────────────────────────────────
s = blank_slide(prs)
bg(s, WHITE)
add_green_bar(s, t=Inches(0), thickness=Inches(0.08))
section_header_block(s, "06", "비용 분석")

# ── 왼쪽 영역: 비용 테이블 ───────────────────
txt(s, "실제 발생 비용  (14개 카드 기준)",
    Inches(0.4), Inches(1.38), Inches(7.5), Inches(0.4),
    size=14, bold=True, color=NAVY)

# 테이블 헤더
TBL_L = Inches(0.4)
TBL_TOP = Inches(1.88)
COL_WS  = [Inches(2.4), Inches(3.4), Inches(1.5)]
COL_LS  = [TBL_L, TBL_L+Inches(2.4), TBL_L+Inches(5.8)]
HDR_TXT = ["모델", "내용", "비용(KRW)"]

for ci, (hl, cw, cl) in enumerate(zip(HDR_TXT, COL_WS, COL_LS)):
    box(s, cl, TBL_TOP, cw, Inches(0.4), bg_color=NAVY)
    txt(s, hl, cl+Inches(0.12), TBL_TOP+Inches(0.06), cw-Inches(0.1), Inches(0.3),
        size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

c_rows = [
    ("OpenAI GPT-image-1", "표지 이미지 (62회, 재생성 포함)", "약 ₩9,800"),
    ("Google Gemini 2.5",  "섹션 이미지 생성",               "₩6,382"),
]
for ri, row in enumerate(c_rows):
    rt = TBL_TOP + Inches(0.4) + ri * Inches(0.45)
    bg_c = LIGHT_GRAY if ri == 0 else WHITE
    for ci, (cell, cw, cl) in enumerate(zip(row, COL_WS, COL_LS)):
        box(s, cl, rt, cw, Inches(0.45), bg_color=bg_c)
        al = PP_ALIGN.RIGHT if ci == 2 else PP_ALIGN.LEFT
        txt(s, cell, cl+Inches(0.12), rt+Inches(0.09), cw-Inches(0.18), Inches(0.3),
            size=12, color=DARK_GRAY, align=al)

# 합계 행
sum_t = TBL_TOP + Inches(1.3)
box(s, TBL_L, sum_t, Inches(7.3), Inches(0.45), bg_color=GREEN)
txt(s, "합    계", TBL_L+Inches(0.12), sum_t+Inches(0.09), Inches(4.0), Inches(0.3),
    size=13, bold=True, color=WHITE)
txt(s, "약 ₩16,200", TBL_L+Inches(4.9), sum_t+Inches(0.09), Inches(2.3), Inches(0.3),
    size=13, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

# 26개 예상 비용
txt(s, "26개 전체 예상 비용",
    Inches(0.4), Inches(3.65), Inches(7.5), Inches(0.4),
    size=14, bold=True, color=NAVY)

p_rows = [
    ("기 발생 비용 (14개 완성)",       "약 ₩16,200"),
    ("추가 예상 (12개 · 스타일 확정)", "약 ₩8,400"),
]
for ri, (label, val) in enumerate(p_rows):
    rt = Inches(4.15) + ri * Inches(0.45)
    box(s, TBL_L, rt, Inches(7.3), Inches(0.45),
        bg_color=LIGHT_GRAY if ri == 0 else WHITE)
    txt(s, label, TBL_L+Inches(0.12), rt+Inches(0.09), Inches(5.2), Inches(0.3),
        size=12, color=DARK_GRAY)
    txt(s, val, TBL_L+Inches(4.9), rt+Inches(0.09), Inches(2.3), Inches(0.3),
        size=12, color=DARK_GRAY, align=PP_ALIGN.RIGHT)

# 총합 강조
box(s, TBL_L, Inches(5.05), Inches(7.3), Inches(0.5), bg_color=NAVY)
txt(s, "총 예상 비용", TBL_L+Inches(0.12), Inches(5.12), Inches(4.5), Inches(0.35),
    size=13, bold=True, color=WHITE)
txt(s, "약 ₩24,000 ~ 30,000", TBL_L+Inches(3.9), Inches(5.12), Inches(3.3), Inches(0.35),
    size=13, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)

# ── 오른쪽 영역: 기존 방식 대비 효율 ────────
box(s, Inches(8.3), Inches(1.38), Inches(4.7), Inches(0.42), bg_color=NAVY)
txt(s, "기존 방식 대비 효율 비교",
    Inches(8.42), Inches(1.44), Inches(4.5), Inches(0.32),
    size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

eff_rows = [
    ("제작 기간",  "9 ~ 11일",         "1일 이내"),
    ("수정 대응",  "디자이너 일정 의존", "즉시 가능"),
    ("총 비용",    "수십만 원 이상",    "₩24,000~30,000"),
    ("파일 크기",  "원본 대용량",       "WebP 96% 감소"),
]

# 헤더 행
for ci, (label, cw, cl) in enumerate(zip(
        ["구분", "기존", "신규"],
        [Inches(1.3), Inches(1.7), Inches(1.6)],
        [Inches(8.3), Inches(9.6), Inches(11.3)])):
    bg_c = RGBColor(0x22,0x3A,0x5E) if ci > 0 else RGBColor(0x22,0x3A,0x5E)
    box(s, cl, Inches(1.88), cw, Inches(0.38), bg_color=RGBColor(0x22,0x3A,0x5E))
    fc = RGBColor(0xFF,0x88,0x88) if ci == 1 else (GREEN if ci == 2 else WHITE)
    txt(s, label, cl+Inches(0.08), Inches(1.92), cw-Inches(0.1), Inches(0.28),
        size=11, bold=True, color=fc, align=PP_ALIGN.CENTER)

for ri, (cat, bef, aft) in enumerate(eff_rows):
    rt = Inches(2.26) + ri * Inches(0.55)
    bg_c = LIGHT_GRAY if ri % 2 == 0 else WHITE
    for ci, (cell, cw, cl) in enumerate(zip(
            [cat, bef, aft],
            [Inches(1.3), Inches(1.7), Inches(1.6)],
            [Inches(8.3), Inches(9.6), Inches(11.3)])):
        box(s, cl, rt, cw, Inches(0.55), bg_color=bg_c)
        fc = RGBColor(0xFF,0x55,0x55) if ci == 1 else (GREEN if ci == 2 else DARK_GRAY)
        bld = ci == 2
        txt(s, cell, cl+Inches(0.08), rt+Inches(0.11), cw-Inches(0.1), Inches(0.35),
            size=11, bold=bld, color=fc, align=PP_ALIGN.CENTER)

# 비고 한 줄
txt(s, "* 테스트·재생성 포함 실비용 기준  |  스타일 확정 후 재생성 대폭 감소",
    Inches(0.4), Inches(6.8), Inches(12.5), Inches(0.38),
    size=10, color=MID_GRAY)


# ──────────────────────────────────────────────
# SLIDE 9 — 향후 계획
# ──────────────────────────────────────────────
s = blank_slide(prs)
bg(s, WHITE)
add_green_bar(s, t=Inches(0), thickness=Inches(0.08))
section_header_block(s, "07", "향후 계획 및 확대 활용 방향")

txt(s, "이번 시스템을 기반으로 다양한 운영 콘텐츠 영역으로 확대 적용을 검토합니다.",
    Inches(0.4), Inches(1.35), Inches(12.5), Inches(0.42),
    size=13, color=MID_GRAY)

# 3개 영역 카드
plan_areas = [
    (
        "01",
        "편집기 활용 영역 확대",
        GREEN,
        [
            ("검진센터 이벤트 팝업 이미지",  "항목별 이벤트 안내 팝업을 편집기로 직접 제작 · 수정"),
            ("운영 콘텐츠 확대 적용",       "다른 건강 주제 카드뉴스, 시즌 캠페인 배너 등"),
            ("다양한 포맷 편집기 제작",      "팝업 · 배너 등 포맷별 맞춤 편집기 구성 검토"),
        ],
    ),
    (
        "02",
        "파이프라인 확장 검토",
        NAVY,
        [
            ("CMS 자동 업로드 연동",         "완성 이미지를 CMS에 직접 등록하는 연동 검토"),
            ("기획서 → 자동화 범위 확대",    "엑셀 기획서 파싱 → 배치 생성 자동화 적용"),
            ("타 콘텐츠 영역 파이프라인 구축", "필요 영역별 전용 파이프라인 구성 단계적 추진"),
        ],
    ),
    (
        "03",
        "효율성 검토 및 고도화",
        RGBColor(0x1A, 0x5C, 0x3A),
        [
            ("활용 부서 · 채널 확대 검토",    "타 부서 콘텐츠 제작 수요 파악 및 적용 가능성 검토"),
            ("제작 비용 · 시간 절감 효과 측정", "도입 전후 KPI 비교를 통한 효과 정량화"),
            ("편집기 기능 고도화",            "섹션 순서 변경, 텍스트 강조 스타일 등 UX 개선"),
        ],
    ),
]

for i, (no, area_title, hdr_c, items) in enumerate(plan_areas):
    left = Inches(0.35) + i * Inches(4.3)
    cw = Inches(4.05)

    # 헤더 바
    box(s, left, Inches(1.88), cw, Inches(0.5), bg_color=hdr_c)
    txt(s, no, left+Inches(0.15), Inches(1.93), Inches(0.4), Inches(0.38),
        size=13, bold=True, color=RGBColor(0xFF,0xFF,0xFF) if hdr_c != GREEN else NAVY,
        align=PP_ALIGN.CENTER)
    txt(s, area_title, left+Inches(0.55), Inches(1.94), cw-Inches(0.65), Inches(0.38),
        size=14, bold=True, color=WHITE)

    # 세부 항목
    for j, (item_title, item_desc) in enumerate(items):
        it = Inches(2.45) + j * Inches(1.5)
        box(s, left, it, cw, Inches(1.42),
            bg_color=LIGHT_GRAY if j % 2 == 0 else WHITE,
            border_color=RGBColor(0xDD,0xDD,0xDD), border_pt=0.5)
        # 왼쪽 포인트 바
        box(s, left, it, Inches(0.06), Inches(1.42), bg_color=hdr_c)
        txt(s, item_title,
            left+Inches(0.15), it+Inches(0.18), cw-Inches(0.22), Inches(0.38),
            size=12, bold=True, color=NAVY)
        txt(s, item_desc,
            left+Inches(0.15), it+Inches(0.62), cw-Inches(0.22), Inches(0.65),
            size=11, color=DARK_GRAY)

# 하단 메시지
box(s, Inches(0.35), Inches(7.02), Inches(12.6), Inches(0.38), bg_color=GREEN_LIGHT)
txt(s, "현재 간 건강 콘텐츠에서 검증된 파이프라인을 기반으로 — 필요 영역에 단계적으로 확대 적용하여 제작 효율성을 높입니다.",
    Inches(0.5), Inches(7.06), Inches(12.3), Inches(0.3),
    size=11, color=NAVY, align=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────
# SLIDE 10 — 마무리
# ──────────────────────────────────────────────
s = blank_slide(prs)
bg(s, NAVY)
box(s, 0, 0, Inches(0.18), H, bg_color=GREEN)

circ_e = s.shapes.add_shape(9, Inches(10.5), Inches(5.2), Inches(3.5), Inches(3.5))
circ_e.fill.solid(); circ_e.fill.fore_color.rgb = RGBColor(0x14, 0x30, 0x55)
circ_e.line.fill.background()

txt(s, "감사합니다",
    Inches(0.8), Inches(2.2), Inches(10), Inches(1.5),
    size=52, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

txt(s, "GC Care 디지털콘텐츠팀  |  2026년 5월",
    Inches(0.8), Inches(3.9), Inches(9), Inches(0.55),
    size=16, color=RGBColor(0x88, 0xAA, 0xCC))

summary = [
    "✓  AI 자동화로 182컷 제작 기간 9일 → 1일 이내",
    "✓  웹 편집기로 담당자 직접 수정 — 디자이너 반복 요청 제거",
    "✓  26개 카드뉴스 약 ₩24,000~30,000 비용으로 완성",
]
for i, line in enumerate(summary):
    txt(s, line, Inches(0.8), Inches(4.65) + i * Inches(0.52), Inches(10), Inches(0.48),
        size=14, color=RGBColor(0xCC, 0xEE, 0xFF))


# ── 저장 ──────────────────────────────────────
out = "카드뉴스_제작시스템_도입보고.pptx"
prs.save(out)
print("저장 완료: " + out)
